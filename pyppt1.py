from pptx import Presentation
from collection import Container
# List of downloaded image filenames
image_filenames = ['image1 (1).jpg', 'image (12).jpg', 'image (13).jpg','image (14).jpg']

# Create a new PowerPoint presentation
presentation = Presentation()

# Loop through the image filenames and add slides with images
for filename in image_filenames:
    # Create a slide with a blank layout
    slide_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(slide_layout)

    # Define the width and height of the image on the slide (adjust as needed)
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height
    image_width = slide_width
    image_height = slide_height

    # Calculate the position to center the image on the slide
    left = (slide_width - image_width) / 2
    top = (slide_height - image_height) / 2

    # Add the image to the slide
    slide.shapes.add_picture(filename, left, top, image_width, image_height)

# Save the PowerPoint presentation
presentation.save('image_presentation.pptx')
